"""Deliverable generators: DOCX, XLSX, PPTX, PDF.

All produce real files saved to the on-premise deliverable directory.
"""
from __future__ import annotations

import io
import uuid
from pathlib import Path

from app.core.config import get_settings


def _new_name(category: str, base: str) -> tuple[str, Path]:
    settings = get_settings()
    stem = base.replace(" ", "_")[:40] or "deliverable"
    safe = "".join(c for c in stem if c.isalnum() or c in "_-")
    fname = f"{safe}_{uuid.uuid4().hex[:6]}.{category}"
    path = settings.deliverable_dir / fname
    path.parent.mkdir(parents=True, exist_ok=True)
    return fname, path


def generate_docx(title: str, sections: list[tuple[str, str]], output: str | None = None) -> tuple[str, Path]:
    from docx import Document
    from docx.shared import Pt
    doc = Document()
    doc.add_heading(title, level=0)
    for head, body in sections:
        doc.add_heading(head, level=1)
        for para in body.split("\n"):
            p = doc.add_paragraph(para)
            p.paragraph_format.space_after = Pt(6)
    if output is None:
        fname, path = _new_name("docx", title)
    else:
        fname, path = output, get_settings().deliverable_dir / output
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    return fname, path


def generate_xlsx(rows: list[list], headers: list[str] | None, sheet: str, output: str | None = None) -> tuple[str, Path]:
    import openpyxl
    from openpyxl.styles import Font, PatternFill
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet[:31] or "Sheet1"
    if headers:
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="0EA5E9")
    for r in rows:
        ws.append(["" if v is None else v for v in r])
    for col in ws.columns:
        width = max((len(str(c.value)) if c.value else 8) for c in col) + 2
        ws.column_dimensions[col[0].column_letter].width = min(width, 32)
    if output is None:
        fname, path = _new_name("xlsx", sheet)
    else:
        fname, path = output, get_settings().deliverable_dir / output
    path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(path))
    return fname, path


def generate_pptx(title: str, slides: list[tuple[str, str]], output: str | None = None) -> tuple[str, Path]:
    from pptx import Presentation
    prs = Presentation()
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = title
    if slides:
        s0.placeholders[1].text = slides[0][0]
    for head, body in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = head
        slide.placeholders[1].text = body
    if output is None:
        fname, path = _new_name("pptx", title)
    else:
        fname, path = output, get_settings().deliverable_dir / output
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return fname, path


def generate_pdf(text: str, output: str | None = None) -> tuple[str, Path]:
    import fitz
    doc = fitz.open()
    margin = 50
    fontsize = 10
    line_h = 14
    usable_w = doc[0].rect.width - 2 * margin
    page = None
    y = None

    def fresh_page():
        nonlocal page, y
        page = doc.new_page()
        y = margin

    fresh_page()
    for para in text.split("\n"):
        words = para.split(" ")
        line = ""
        for w in words:
            trial = (line + " " + w).strip()
            if len(trial) * 6.2 > usable_w:  # approx char width
                if y > page.rect.height - margin:
                    fresh_page()
                page.insert_text((margin, y), line, fontsize=fontsize, fontname="helv")
                y += line_h
                line = w
            else:
                line = trial
        if line:
            if y > page.rect.height - margin:
                fresh_page()
            page.insert_text((margin, y), line, fontsize=fontsize, fontname="helv")
            y += line_h

    if output is None:
        fname, path = _new_name("pdf", "analysis")
    else:
        fname, path = output, get_settings().deliverable_dir / output
    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))
    return fname, path


def save_code(code: str, language: str, output: str | None = None) -> tuple[str, Path]:
    ext = {"python": "py", "javascript": "js", "shell": "sh"}.get(language, "txt")
    if output is None:
        fname, path = _new_name(ext, "generated_code")
    else:
        fname, path = output, get_settings().deliverable_dir / output
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(code)
    return fname, path
